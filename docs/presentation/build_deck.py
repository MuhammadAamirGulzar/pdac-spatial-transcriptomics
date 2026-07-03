"""
Builds the results presentation, matching the team's existing deck design
(Liter font, blue 2C5F8A / orange E07B54, blue header bars, PART 0X dividers).
Every result slide shows a figure (the proof) with a short takeaway and speaker notes.
Output: docs/meetings/Meeting - Results and Findings (PDAC ST + Foundation Model).pptx
"""
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from PIL import Image

ROOT = os.path.dirname(os.path.abspath(__file__)); PROJ = os.path.dirname(os.path.dirname(ROOT))
FIG = os.path.join(PROJ, "Outputs", "presentation_figures")
SL = os.path.join(FIG, "slides")
CACHE = os.path.join(ROOT, "_tissue_cache")

BLUE = RGBColor(0x2C, 0x5F, 0x8A); ORANGE = RGBColor(0xE0, 0x7B, 0x54)
LIGHT = RGBColor(0xD9, 0xE8, 0xF5); GREY = RGBColor(0x8D, 0x8D, 0x8D)
DARK = RGBColor(0x33, 0x33, 0x33); WHITE = RGBColor(0xFF, 0xFF, 0xFF)
NAVY = RGBColor(0x1E, 0x4A, 0x6E)
FONT = "Liter"
W, H = 17.7778, 10.0

prs = Presentation(); prs.slide_width = Inches(W); prs.slide_height = Inches(H)
BLANK = prs.slide_layouts[6]

_page = [0]                       # auto slide counter (every slide occupies a number)
def slide():
    _page[0] += 1
    return prs.slides.add_slide(BLANK)

def rect(s, l, t, w, h, color, line=False):
    sp = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(l), Inches(t), Inches(w), Inches(h))
    sp.fill.solid(); sp.fill.fore_color.rgb = color
    if not line: sp.line.fill.background()
    sp.shadow.inherit = False
    return sp

def txt(s, text, l, t, w, h, size, color=DARK, bold=False, align=PP_ALIGN.LEFT,
        anchor=MSO_ANCHOR.TOP, font=FONT, italic=False, line_spacing=None):
    tb = s.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = True; tf.vertical_anchor = anchor
    lines = text.split("\n")
    for i, ln in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        if line_spacing: p.line_spacing = line_spacing
        r = p.add_run(); r.text = ln
        f = r.font; f.name = font; f.size = Pt(size); f.bold = bold; f.italic = italic
        f.color.rgb = color
    return tb

def footer(s):
    rect(s, 0, 9.94, W, 0.06, ORANGE)

def header(s, title, page=None):
    rect(s, 0, 0, W, 0.83, BLUE)
    rect(s, 0.33, 0.19, 0.44, 0.44, ORANGE)   # small accent square
    txt(s, title, 0.97, 0.0, 14.5, 0.83, 26, WHITE, bold=True, anchor=MSO_ANCHOR.MIDDLE)
    txt(s, f"{_page[0]:02d}", 16.7, 0.0, 0.9, 0.83, 16, LIGHT, anchor=MSO_ANCHOR.MIDDLE, align=PP_ALIGN.CENTER)
    footer(s)

def fitted(s, path, bl, bt, bw, bh):
    iw, ih = Image.open(path).size; ar = iw / ih; bar = bw / bh
    if ar > bar: w = bw; h = bw / ar
    else: h = bh; w = bh * ar
    s.shapes.add_picture(path, Inches(bl + (bw - w) / 2), Inches(bt + (bh - h) / 2), width=Inches(w))

def notes(s, text):
    s.notes_slide.notes_text_frame.text = text

# ---------------------------------------------------------------- 1 TITLE
def title_slide():
    s = slide()
    rect(s, 11.85, 0, W - 11.85, H, NAVY)
    tpath = os.path.join(CACHE, "IU_PDA_T3_he.png")
    if os.path.exists(tpath): fitted(s, tpath, 12.15, 2.7, 5.3, 4.6)
    txt(s, "TISSUE TO MOLECULES TO PREDICTION", 12.15, 7.5, 5.3, 0.4, 12, LIGHT, bold=True, align=PP_ALIGN.CENTER)
    txt(s, "SPATIAL TRANSCRIPTOMICS RESEARCH", 1.11, 1.5, 9.7, 0.4, 15, GREY, bold=True)
    txt(s, "Spatial Transcriptomics and Multimodal\nMetastatic Biomarker Discovery",
        1.11, 2.1, 10.2, 2.0, 44, BLUE, bold=True, line_spacing=1.05)
    txt(s, "Results and What We Learned", 1.11, 4.8, 10, 0.6, 22, ORANGE, bold=True)
    rect(s, 1.13, 5.6, 2.4, 0.04, ORANGE)
    txt(s, "Research Advisors:  Dr. Ashiq Masood and Dr. Mohsin Bilal\nDate:  29 June 2026",
        1.11, 5.9, 9.7, 1.2, 17, GREY, line_spacing=1.25)
    footer(s)
    notes(s, "This is a results update. The goal of the project is to flag which tumour regions are prone "
              "to spread to the liver, using a routine H and E slide at the time of prediction. I will walk "
              "from the question, to how we built a trustworthy signal, to what a plain slide can and cannot "
              "tell us, and finish with an honest verdict and next steps.")

# ---------------------------------------------------------------- 2 OVERVIEW
def overview_slide():
    s = slide()
    txt(s, "What this talk covers", 1.11, 0.56, 10, 0.7, 30, BLUE, bold=True)
    rect(s, 1.11, 1.33, 15.56, 0.02, LIGHT)
    items = [("1", "The question and the data", "What we predict and the six tissues"),
             ("2", "A signal we can trust", "The confound, the leaving program, one worked example"),
             ("3", "What a plain H&E can do", "How far the slide alone gets us"),
             ("4", "What we learned and next", "Honest scorecard and the way forward")]
    xw = 3.75; x0 = 1.0
    for i, (n, t, d) in enumerate(items):
        x = x0 + i * xw
        c = rect(s, x + 0.3, 2.5, 0.56, 0.56, BLUE);
        txt(s, n, x + 0.3, 2.5, 0.56, 0.56, 20, WHITE, bold=True, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        txt(s, t, x, 3.3, xw - 0.2, 0.8, 19, BLUE, bold=True)
        txt(s, d, x, 4.0, xw - 0.2, 1.2, 15, GREY)
    txt(s, "Short version:  the biology we found is real and checks out against outside data. "
           "A plain slide today reads how much tumour is present, and not yet the finer movement state.",
        1.11, 5.7, 15.4, 1.4, 19, DARK, line_spacing=1.2)
    footer(s)
    notes(s, "Four parts. First the question and the six samples. Second, how we built a signal we can trust, "
              "including the one confound that could have fooled us, and a single spot worked through end to end. "
              "Third, how far a plain slide gets us on new patients. Fourth, the honest scorecard and what would push it further.")

# ---------------------------------------------------------------- SECTION DIVIDER
def divider(part, sub, big, subtitle, note):
    s = slide()
    rect(s, 0, 0, W, 3.0, BLUE)
    txt(s, f"PART  {part}", 1.11, 0.83, 8, 0.6, 24, LIGHT)
    txt(s, sub, 1.11, 1.55, 12, 0.5, 18, LIGHT)
    rect(s, 1.11, 3.75, 0.06, 1.67, ORANGE)
    txt(s, big, 1.44, 3.6, 15, 1.9, 46, BLUE, bold=True)
    txt(s, subtitle, 1.44, 5.8, 13.5, 1.0, 20, GREY)
    footer(s)
    notes(s, note)

# ---------------------------------------------------------------- FIGURE SLIDE
def fig_slide(title, page, fig_path, takeaway, note, img_h=7.7):
    s = slide(); header(s, title, page)
    fitted(s, fig_path, 0.35, 1.0, 17.1, img_h)
    if takeaway:
        rect(s, 0.0, 9.18, W, 0.62, LIGHT)
        txt(s, takeaway, 0.6, 9.18, 16.6, 0.62, 16, DARK, bold=True, anchor=MSO_ANCHOR.MIDDLE)
    footer(s)
    notes(s, note)

# ---------------------------------------------------------------- NEXT STEPS (cards)
def next_slide(page):
    s = slide(); header(s, "What would push this further", page)
    cards = [("01", "Pathologist annotation", "Mark invasive fronts on a few slides so we can test the signal against expert ground truth, not just transcriptomics.", BLUE),
             ("02", "Higher resolution", "Move from 55 micron spots to Visium HD or Xenium, so a spot is closer to a single cell and fine signal survives.", ORANGE),
             ("03", "More matched patients", "Add patients that have both a primary tumour and its own liver metastasis, so the link can be learned, not guessed.", BLUE)]
    cw = 5.2; x0 = 0.83
    for i, (n, t, d, c) in enumerate(cards):
        x = x0 + i * (cw + 0.25)
        rect(s, x, 1.3, cw, 5.6, WHITE, line=True)
        rect(s, x, 1.3, cw, 0.08, c)
        rect(s, x + 0.28, 1.62, 1.0, 0.5, c)
        txt(s, "STEP " + n, x + 0.28, 1.62, 1.0, 0.5, 12, WHITE, bold=True, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        txt(s, t, x + 0.28, 2.4, cw - 0.5, 0.8, 21, c, bold=True)
        txt(s, d, x + 0.28, 3.4, cw - 0.5, 3.0, 16, DARK, line_spacing=1.2)
    rect(s, 0.83, 7.4, 16.1, 1.6, LIGHT)
    txt(s, "The biology is already real and validated.  These steps are about the read out from the slide, "
           "not about the signal existing.  All three can run in parallel.",
        1.2, 7.4, 15.4, 1.6, 18, DARK, bold=True, anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.2)
    footer(s)
    notes(s, "None of these change the finding that the program is real. They are about making the slide able to "
              "read the finer state. Annotation gives us expert ground truth. Higher resolution gives cleaner signal "
              "per spot. More matched patients let the primary to metastasis link be learned directly.")

# ---------------------------------------------------------------- CLOSING
def closing_slide(page):
    s = slide()
    rect(s, 0, 0, W, H, BLUE)
    txt(s, "Thank you", 1.5, 3.1, 15, 1.2, 48, WHITE, bold=True)
    rect(s, 1.55, 4.5, 2.6, 0.05, ORANGE)
    txt(s, "The intra-primary EMT leaving program is genuine and externally validated.\n"
           "A routine H&E slide currently reads how much tumour is present, and not yet the finer\n"
           "metastatic state on new patients.  Your feedback on the invasive front annotation would help most.",
        1.5, 5.0, 15, 2.2, 20, LIGHT, line_spacing=1.3)
    txt(s, "Spatial Transcriptomics Research   |   29 June 2026", 1.5, 8.6, 12, 0.5, 15, LIGHT)
    footer(s)
    notes(s, "Happy to take questions. The one ask that would move us fastest is a small set of slides with the "
              "invasive fronts marked, so we can test the score against your read rather than against transcriptomics alone.")

# ================================================================ BUILD
title_slide()
overview_slide()

divider("01", "The clinical question and the cohort", "The question and the data",
        "What we predict, and the six tissues it is built on",
        "We want to read metastasis risk from a routine slide. Spatial transcriptomics is the teacher during "
        "training only. At prediction time we want only the H and E slide.")
fig_slide("The question and the data", 4, os.path.join(SL, "fig1_cohort_overview.png"),
          "Goal: flag metastasis prone tumour regions from a routine H&E slide.  Transcriptomics teaches the model during training only.",
          "Six tissues. Four primary pancreatic tumours and two liver metastases. Each spot has three matched readings, "
          "the H and E image, the gene expression, and the cell makeup. Only patient 11 has both a primary and its own "
          "metastasis. The constraint from the clinic is that at prediction time we use the slide alone.")

divider("02", "Building a signal we can trust", "A signal we can trust",
        "The confound we avoided, the program we measure, and one spot worked through end to end",
        "Before trusting any score we had to rule out the obvious trap, then define a signal inside the primary "
        "tumour, then show one spot end to end so it is not a black box.")
fig_slide("Liver contamination: how much, and why it fakes a signal", 6, os.path.join(SL, "figLsum_howmuch_removable.png"),
          "Metastasis spots average ~15% liver cells vs ~0% in primary.  That liver content IS the naive 'metastasis' axis (corr 0.40), and restricting to tumour-only spots removes it (0.07).",
          "Two panels, one message. Left, how much: liver metastasis spots are about fifteen percent liver cells, the thirteen thousand "
          "five hundred primary spots almost none. Right, why it matters: the naive metastasis direction correlates zero point four with "
          "that liver content across all eighteen thousand spots, and restricting to the four thousand eight hundred tumour-dominated "
          "spots collapses it to zero point zero seven. Most of the apparent metastasis signal was simply the liver the tumour sits in.")
fig_slide("Where the liver cells are — metastasis HM11", 7, os.path.join(SL, "figLwhere_HM11.png"),
          "On the real HM11 slide, the liver-rich zone zooms to classic liver parenchyma (cords + sinusoids); the tumour zone is dense desmoplasia with no liver.",
          "Here is the contamination on the real HM11 whole slide, spots coloured by liver-cell content. Zoom into the bright pocket, cyan "
          "box, and the tissue has the cords and sinusoids of normal liver. Zoom into the tumour zone, green box, and it is dense malignant "
          "and fibrous tissue with no liver. The rings are left unfilled so you can see the underlying H and E. The tumour is physically "
          "embedded in the liver it landed in.")
fig_slide("Where the liver cells are — metastasis HM13", 8, os.path.join(SL, "figLwhere_HM13.png"),
          "The second metastasis tells the same story, with even more liver content (mean ~21%): a clear liver zone and a clear tumour zone on the real tissue.",
          "The second metastasis, HM13, tells the same story, and it actually has more liver, about twenty one percent on average. Same "
          "layout, the liver-rich zone zooms to liver parenchyma and the tumour zone to malignant tissue. Two independent metastasis "
          "slides, the same liver contamination, which is exactly the confound we must avoid.")
fig_slide("The primary tumour has essentially no liver", 9, os.path.join(SL, "figLwhere_T1.png"),
          "By contrast, the primary tumour is all tumour/pancreas — the maximum liver-cell fraction anywhere is ~22%, and typical tissue is malignant desmoplasia.  The confound is metastasis-only.",
          "By contrast, the primary tumour. Every spot is dark, near zero liver content. Even the single most liver-like spot is barely "
          "any, and the zoomed tissue is malignant and desmoplastic, not liver. This is the clean control. The liver confound lives only "
          "in the metastasis slides, which is precisely why we define the target inside the primary tumour.")
fig_slide("Proof it is really liver — the cells and the genes agree", 10, os.path.join(SL, "figL2_liver_proof.png"),
          "The same spots carry liver cells AND liver genes (ALB, TTR, APOA1).  Cell deconvolution and gene expression agree at r = 0.93.",
          "And the proof it is liver and not a deconvolution quirk. Group means over all liver-rich metastasis spots versus all tumour "
          "spots: the liver spots strongly express the classic hepatocyte genes, albumin, transthyretin, APOA1, while tumour spots show "
          "essentially none. Across every HM11 spot the cell-composition liver fraction and the liver-gene score agree almost perfectly, "
          "correlation zero point nine three. The liver contamination is real biology.")
fig_slide("Two metastases do not resemble each other", 11, os.path.join(SL, "figL3_generalization.png"),
          "Hold out one sample and predict it from the rest: primary tumours generalise (93–100%), but neither liver metastasis is recovered from the other (0%).",
          "One more problem, patient number. Leave one sample out and ask whether the model recognises its tumour type from the others. "
          "Primary tumours generalise, ninety three to one hundred percent. But each liver metastasis is not recovered from the other at "
          "all, zero percent. With only two metastasis patients there is no transferable metastasis signature to learn. Together with the "
          "liver confound, this is why we do not model the metastasis tissue directly.")
fig_slide("A movement and invasion program inside the tumour", 7, os.path.join(SL, "fig3_stage1a_leaving_maps.png"),
          "High score regions cluster together in every tumour, like invasive fronts, not random noise.",
          "This is the leaving program. For each tumour, the H and E is on the left and the score is on the right, on the same tissue. "
          "Red marks regions expressing the movement and invasion genes. The high regions form coherent zones in every tumour, "
          "which is what you would expect for invasive fronts, and not what random noise looks like.")
fig_slide("It is real, known EMT biology", 8, os.path.join(SL, "fig4_stage1a_program_biology.png"),
          "Built from 19 textbook EMT genes.  An independent paper's own EMT score agrees with ours.",
          "The score is the average of nineteen well known epithelial to mesenchymal transition and invasion genes, the genes "
          "tumours use to detach and move. On the right, computed completely separately, our score lines up with the source "
          "paper's own EMT and stroma signatures, and not with unrelated programs. So this is genuine dissemination biology.")
fig_slide("Which genes actually carry the signal", 9, os.path.join(SL, "figB_leaving_gene_importance.png"),
          "The confound-free signal is driven by secreted matrix and protease genes (SERPINE1, S100A4, collagens), not the textbook EMT switches.",
          "Not all nineteen genes contribute equally. Averaged over all four tumours, the confound free signal is carried by the "
          "secreted matrix and protease genes, SERPINE1, S100A4, collagen five, laminin, fibronectin, and not by the classic "
          "transcription factor switches like SNAI1, which are low count and noisy at this fifty five micron resolution. On the "
          "right, collagens track how much tumour is present, while tenascin and laminin track the finer intrinsic invasive state. "
          "Data driven across all seventeen thousand genes, the top hits are exactly our panel genes, which independently confirms the choice.")
fig_slide("The confound and the fix, on three real spots", 10, os.path.join(SL, "figA_leaving_worked_spots.png"),
          "On the real slide, each spot's own H&E is shown so you can judge it by eye.  A dense spot can score high just from more tumour; removing that shows C is bulky, not invasive, while A and B are genuinely leaving.",
          "Here is the one confound that could fool us, made concrete on three real spots from tumour four, shown on the actual whole "
          "slide with each spot's own H&E patch so a pathologist can check every call by eye. The raw score climbs with tumour amount, "
          "the red line. Spot C has eighty two percent tumour but is a bulky sheet, not invasive. Spots A and B have less tumour but "
          "fibrous, infiltrative edges and strongly express the program. After we remove the tumour amount effect, A and B stay high "
          "while C flips negative. That corrected residual is the target we actually model, so we measure invasiveness, not density.")
fig_slide("Reading one result end to end", 11, os.path.join(FIG, "fig_worked_example.png"),
          "Location, then H&E you can judge by eye, then genes, then score, then outside confirmation.  Not a black box.",
          "Here is one result made fully traceable. Two real tumour spots, one high and one low. You can see where each sits, "
          "then its actual H and E, where the high spot looks pink and fibrous and the low one densely cellular. On the right, "
          "across all tumour spots, the high score spots switch the EMT genes on while the epithelial genes stay flat. And each "
          "spot also carries an outside paper's EMT score that agrees. A pathologist can check every step.")

divider("03", "From transcriptomics to a plain slide", "What a plain H&E can do",
        "How far the slide alone gets us on patients it never saw",
        "Now the real test. Can the slide alone reproduce the signal on new patients, which is the clinical promise.")
fig_slide("What an H&E slide can and cannot tell us", 13, os.path.join(SL, "fig5_stages23_ceiling.png"),
          "The slide reliably reads how much tumour is present.  The finer movement signal is not yet readable across patients.",
          "On new patients the slide reproduces the raw score moderately. But once we remove the simple how much tumour is here "
          "effect, the finer movement signal drops to the noise floor. A trivial tumour amount predictor actually beats our model. "
          "And adding the transcriptomic bridge did not beat the foundation model on its own. So today the slide mainly senses "
          "tumour amount, not the fine metastatic state. This is the honest ceiling, and it is better to know it now.")
fig_slide("Predicted vs real on the tissue — tumour T4", 14, os.path.join(SL, "figP_T4.png"),
          "Whole slide, actual score, and H&E prediction.  Zooming a green (agree) and a red (disagree) region shows the prediction tracks broad zonation but misses finer structure.",
          "The same story on the real tissue, one tumour at a time. T4: the whole slide, the actual transcriptomic score, and the score "
          "predicted from H and E alone. The green box marks a region where the two agree, the red box a region where they disagree. "
          "Zoom in and the thin rings, actual on the left, predicted on the right, match in the agree region and flip colour in the "
          "disagree region. The prediction follows the broad high and low zones but misses the finer structure.")
fig_slide("Predicted vs real on the tissue — tumour T1", 15, os.path.join(SL, "figP_T1.png"),
          "The second tumour tells the same story: broad zonation is captured from H&E, fine structure is not — the agree/disagree zooms make it concrete.",
          "The second tumour, T1, tells the same story. Broad zonation is recovered from the slide, but where we zoom into a disagree "
          "region the predicted rings no longer match the actual ones. This is the visual version of the ceiling, made checkable spot "
          "by spot on the real tissue.")

divider("04", "Where we stand", "What we learned and next",
        "An honest scorecard, and the steps that would push this further",
        "To close, an honest scorecard against pre stated tests, then the steps that would move it forward.")
fig_slide("Scorecard and honest verdict", 16, os.path.join(SL, "fig7_stage4_scorecard.png"),
          "The program is real and validated.  The slide today reads tumour amount, not the confound free state, on new patients.",
          "Six checks we set in advance. The target is real EMT and spatially organised, both yes. The slide reads only the "
          "tumour amount part, it does not beat a trivial tumour amount model, the transcriptomic bridge did not help, and the "
          "primary to metastasis link in patient eleven did not hold. This is a clean, confound checked result. The biology is "
          "real, the slide read out is amount limited at this resolution.")
next_slide(17)
closing_slide(18)

out = os.path.join(ROOT, "Meeting - Results and Findings (PDAC ST + Foundation Model).pptx")
prs.save(out)
print("saved:", out, "| slides:", len(prs.slides.__iter__.__self__._sldIdLst))
