"""
Builds the updated results deck for the clinical review meeting.

Audience is a clinician and pathologist, not a machine learning audience. The
rules followed here:
  - one idea per slide, stated as a sentence in the header, not as a label
  - no model names, no statistics vocabulary on the slide face
  - every claim shown as a picture, with the control that makes it believable
    given its own slide rather than buried in a footnote
  - the negative results get their own slide, before the implications

Speaker notes carry the detail, so the slide can stay almost empty.

Matches the team's existing deck design (Liter font, blue 2C5F8A / orange E07B54,
blue header bars, PART 0X dividers).

Run:
    python docs/presentation/build_deck_site.py
Output:
    docs/presentation/Meeting - Where a tumour spreads changes what it becomes.pptx
"""

import os

from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

ROOT = os.path.dirname(os.path.abspath(__file__))
PROJ = os.path.dirname(os.path.dirname(ROOT))
FIG = os.path.join(PROJ, "Outputs", "presentation_figures")
SS = os.path.join(FIG, "slides_site")
R = os.path.join(PROJ, "Outputs")
CACHE = os.path.join(ROOT, "_tissue_cache")
DEST = os.path.join(ROOT, "Meeting - Where a tumour spreads changes what it becomes.pptx")

BLUE = RGBColor(0x2C, 0x5F, 0x8A); ORANGE = RGBColor(0xE0, 0x7B, 0x54)
LIGHT = RGBColor(0xD9, 0xE8, 0xF5); GREY = RGBColor(0x8D, 0x8D, 0x8D)
DARK = RGBColor(0x33, 0x33, 0x33); WHITE = RGBColor(0xFF, 0xFF, 0xFF)
NAVY = RGBColor(0x1E, 0x4A, 0x6E); GREEN = RGBColor(0x1B, 0xAF, 0x7A)
FONT = "Liter"
W, H = 17.7778, 10.0

prs = Presentation(); prs.slide_width = Inches(W); prs.slide_height = Inches(H)
BLANK = prs.slide_layouts[6]
_page = [0]


def slide():
    _page[0] += 1
    return prs.slides.add_slide(BLANK)


def rect(s, l, t, w, h, color, line=False):
    sp = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(l), Inches(t), Inches(w), Inches(h))
    sp.fill.solid(); sp.fill.fore_color.rgb = color
    if not line:
        sp.line.fill.background()
    sp.shadow.inherit = False
    return sp


def txt(s, text, l, t, w, h, size, color=DARK, bold=False, align=PP_ALIGN.LEFT,
        anchor=MSO_ANCHOR.TOP, font=FONT, italic=False, line_spacing=None):
    tb = s.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = True; tf.vertical_anchor = anchor
    for i, ln in enumerate(text.split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        if line_spacing:
            p.line_spacing = line_spacing
        r = p.add_run(); r.text = ln
        f = r.font; f.name = font; f.size = Pt(size); f.bold = bold; f.italic = italic
        f.color.rgb = color
    return tb


def footer(s):
    rect(s, 0, 9.94, W, 0.06, ORANGE)


def header(s, title):
    rect(s, 0, 0, W, 0.83, BLUE)
    rect(s, 0.33, 0.19, 0.44, 0.44, ORANGE)
    txt(s, title, 0.97, 0.0, 15.3, 0.83, 23, WHITE, bold=True, anchor=MSO_ANCHOR.MIDDLE)
    txt(s, f"{_page[0]:02d}", 16.7, 0.0, 0.9, 0.83, 16, LIGHT,
        anchor=MSO_ANCHOR.MIDDLE, align=PP_ALIGN.CENTER)
    footer(s)


def fitted(s, path, bl, bt, bw, bh):
    iw, ih = Image.open(path).size
    ar, bar = iw / ih, bw / bh
    if ar > bar:
        w, h = bw, bw / ar
    else:
        h, w = bh, bh * ar
    s.shapes.add_picture(path, Inches(bl + (bw - w) / 2), Inches(bt + (bh - h) / 2),
                         width=Inches(w))


def notes(s, text):
    s.notes_slide.notes_text_frame.text = text


def divider(part, sub, big, subtitle, note):
    s = slide()
    rect(s, 0, 0, W, 3.0, BLUE)
    txt(s, f"PART  {part}", 1.11, 0.83, 8, 0.6, 24, LIGHT)
    txt(s, sub, 1.11, 1.55, 13, 0.5, 18, LIGHT)
    rect(s, 1.11, 3.75, 0.06, 1.67, ORANGE)
    txt(s, big, 1.44, 3.6, 15, 1.9, 42, BLUE, bold=True, line_spacing=1.05)
    txt(s, subtitle, 1.44, 5.9, 14.5, 1.2, 20, GREY, line_spacing=1.25)
    footer(s)
    notes(s, note)


def fig_slide(title, fig_path, takeaway, note, img_h=7.6):
    s = slide(); header(s, title)
    fitted(s, fig_path, 0.35, 1.05, 17.1, img_h)
    if takeaway:
        rect(s, 0.0, 9.14, W, 0.66, LIGHT)
        txt(s, takeaway, 0.6, 9.14, 16.6, 0.66, 16, DARK, bold=True, anchor=MSO_ANCHOR.MIDDLE)
    footer(s)
    notes(s, note)


# ================================================================== 1 TITLE
def title_slide():
    s = slide()
    rect(s, 11.85, 0, W - 11.85, H, NAVY)
    tp = os.path.join(CACHE, "IU_PDA_HM11_he.png")
    if os.path.exists(tp):
        fitted(s, tp, 12.15, 2.7, 5.3, 4.6)
    txt(s, "A LIVER METASTASIS, MEASURED POINT BY POINT", 12.15, 7.5, 5.3, 0.5, 11, LIGHT,
        bold=True, align=PP_ALIGN.CENTER)
    txt(s, "SPATIAL TRANSCRIPTOMICS RESEARCH", 1.11, 1.5, 9.7, 0.4, 15, GREY, bold=True)
    txt(s, "Where a tumour spreads\nchanges what it becomes",
        1.11, 2.1, 10.4, 2.2, 44, BLUE, bold=True, line_spacing=1.05)
    txt(s, "Pancreatic cancer in the liver and in a lymph node\nare not the same disease state",
        1.11, 4.75, 10.2, 1.2, 21, ORANGE, bold=True, line_spacing=1.25)
    rect(s, 1.13, 6.05, 2.4, 0.04, ORANGE)
    txt(s, "Research Advisors:  Dr. Ashiq Masood and Dr. Mohsin Bilal\nDate:  14 August 2026",
        1.11, 6.35, 9.7, 1.2, 17, GREY, line_spacing=1.25)
    footer(s)
    notes(s,
          "This is an update on where the project has moved since the last meeting. Last time I showed "
          "that a routine slide can read how much tumour is present but not the finer metastatic state. "
          "Chasing why that was, we ended up asking a more basic question, and that question turned into "
          "the main result. The short version is that a deposit in the liver and a deposit in a lymph "
          "node, from the same patient, are not making the same change. I will show you three separate "
          "datasets that all point the same way, and I will be explicit about what we could not show.")


# ================================================================== 2 OVERVIEW
def overview_slide():
    s = slide()
    txt(s, "What this talk covers", 1.11, 0.56, 10, 0.7, 30, BLUE, bold=True)
    rect(s, 1.11, 1.33, 15.56, 0.02, LIGHT)
    items = [("1", "The question", "What we measure, and the patients it comes from"),
             ("2", "Liver against lymph node", "Are the two deposits doing the same thing?"),
             ("3", "A second, separate group", "Does the same pattern appear in other patients?"),
             ("4", "What it means", "For reading slides, and what we could not show")]
    xw, x0 = 3.75, 1.0
    for i, (n, t, d) in enumerate(items):
        x = x0 + i * xw
        rect(s, x + 0.3, 2.5, 0.56, 0.56, BLUE)
        txt(s, n, x + 0.3, 2.5, 0.56, 0.56, 20, WHITE, bold=True,
            align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        txt(s, t, x, 3.3, xw - 0.2, 0.9, 19, BLUE, bold=True)
        txt(s, d, x, 4.25, xw - 0.2, 1.4, 15, GREY, line_spacing=1.2)
    rect(s, 1.11, 6.1, 15.56, 1.9, LIGHT)
    txt(s, "Short version:  when a pancreatic tumour spreads, only about half of what changes is the "
           "same wherever it lands.\nThe other half is decided by the destination, and the part that "
           "differs is the immune content.",
        1.5, 6.1, 14.8, 1.9, 19, DARK, bold=True, anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.3)
    footer(s)
    notes(s,
          "Four parts. First what we actually measure, because the method is unusual and everything "
          "rests on it. Second the main comparison, liver against lymph node in the same patient. "
          "Third whether it holds up in a completely separate group of patients. Fourth what it means, "
          "including a clear statement of what we could not show. I will keep the statistics off the "
          "slides. Ask me for any number and I have it.")


# ================================================================== COHORTS
def cohorts_slide():
    s = slide(); header(s, "Three separate groups of patients, not one")
    cards = [("01", "The main group", "13 patients, 30 tissue sections\n\n"
              "Primary pancreatic tumour, liver deposits, lymph node deposits, and normal pancreas.\n\n"
              "Four patients gave us both a liver and a lymph node deposit, which is what makes the "
              "key comparison possible.", BLUE),
             ("02", "The check group", "13 patients, 55 tissue sections\n\n"
              "A published cohort from another centre. None of these patients had chemotherapy "
              "before sampling.\n\n"
              "Deposits in liver, lung and peritoneum, so three destinations instead of two.", ORANGE),
             ("03", "The wider test", "156 tissue sections, 7 organs\n\n"
              "A large public collection covering bowel, prostate, kidney, brain, breast, lymph node "
              "and pancreas.\n\n"
              "Used to test whether the idea holds outside pancreatic cancer at all.", GREEN)]
    cw, x0 = 5.2, 0.83
    for i, (n, t, d, c) in enumerate(cards):
        x = x0 + i * (cw + 0.25)
        rect(s, x, 1.3, cw, 6.1, WHITE, line=True)
        rect(s, x, 1.3, cw, 0.08, c)
        rect(s, x + 0.28, 1.62, 1.5, 0.5, c)
        txt(s, "GROUP " + n, x + 0.28, 1.62, 1.5, 0.5, 12, WHITE, bold=True,
            align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        txt(s, t, x + 0.28, 2.35, cw - 0.5, 0.7, 21, c, bold=True)
        txt(s, d, x + 0.28, 3.15, cw - 0.5, 4.0, 14.5, DARK, line_spacing=1.25)
    rect(s, 0.83, 7.75, 16.1, 1.25, LIGHT)
    txt(s, "The second and third groups were never touched while the first was being analysed.  "
           "They are a check, not extra training data.",
        1.2, 7.75, 15.4, 1.25, 18, DARK, bold=True, anchor=MSO_ANCHOR.MIDDLE)
    footer(s)
    notes(s,
          "Three groups, and they do different jobs. The first is where the finding came from. It is "
          "small, thirteen patients, and four of those gave us tissue from two different metastatic "
          "sites, which is rare and is the reason this comparison could be done at all. The second "
          "group is somebody else's published cohort, from another centre, and importantly none of "
          "those patients had chemotherapy before the sample was taken, so treatment cannot explain "
          "anything we see there. It also covers lung and peritoneum, which our group does not. The "
          "third is a large public collection across seven organs and several cancers, used for one "
          "specific test at the end. I want to be clear that we did not look at groups two and three "
          "while developing the finding in group one.")


# ================================================================== NEGATIVES
def negatives_slide():
    s = slide(); header(s, "What we could not show")
    rows = [("The B cell finding rests on one group of patients",
             "The second group had no lymph node samples, so this specific result is untested "
             "elsewhere. It is not contradicted, it is simply not yet checked."),
            ("Four patients is four patients",
             "The paired comparison uses the four people who gave tissue from both sites. The "
             "direction agreed in all four, which is why we trust it, but four is still four."),
            ("We cannot link any of this to how patients did",
             "Clinical details were available for six of the thirteen patients, all at the same "
             "stage, with no survival information. Nothing here is connected to outcome."),
            ("A measuring point is not a single cell",
             "Each point is 55 microns and holds up to ten cells. We restricted to tumour rich "
             "points and repeated everything at stricter thresholds, but we cannot separate a "
             "tumour cell from the cell touching it.")]
    y = 1.35
    for t, d in rows:
        rect(s, 0.9, y, 0.06, 1.55, ORANGE)
        txt(s, t, 1.2, y, 15.4, 0.55, 20, BLUE, bold=True)
        txt(s, d, 1.2, y + 0.62, 15.4, 1.0, 15.5, DARK, line_spacing=1.2)
        y += 1.85
    rect(s, 0.9, 8.75, 16.0, 0.85, LIGHT)
    txt(s, "None of these undo the main result.  They mark where it stops.",
        1.2, 8.75, 15.4, 0.85, 18, DARK, bold=True, anchor=MSO_ANCHOR.MIDDLE)
    footer(s)
    notes(s,
          "I want to put the limits before the implications rather than after, because they change "
          "how much weight you should put on each piece. The B cell result is the one I would most "
          "like to see confirmed, and the honest position is that the second cohort simply had no "
          "lymph nodes in it, so nobody has checked it yet. The paired test involves four people. "
          "With four people you cannot reach a conventional significance threshold no matter how "
          "large the effect is, so what I am relying on is that the direction was the same in every "
          "one of them. And there is no outcome data, so I cannot tell you any of this predicts "
          "survival. If a patient's chart data is available for these cases, that would change what "
          "we can ask next.")


# ================================================================== MEANING
def meaning_slide():
    s = slide(); header(s, "What this means")
    cards = [("For how we think about spread",
              "A metastasis is not one disease state. The organ it lands in shapes roughly half of "
              "what the tumour becomes there. Grouping all metastases together, which is what most "
              "studies do, averages over a real difference.", BLUE),
             ("For immunotherapy thinking",
              "Liver deposits arrived immune poor and lung deposits immune rich, in patients who had "
              "had no chemotherapy. If that holds, the immune environment of a deposit is set "
              "substantially by where it is, not only by the tumour that produced it.", ORANGE),
             ("For AI on pathology slides",
              "A model trained to read tissue in one organ lost more than a third of its accuracy in "
              "another organ, even after allowing for a change of laboratory. A single accuracy "
              "figure across pooled organs will flatter such a model.", GREEN)]
    cw, x0 = 5.2, 0.83
    for i, (t, d, c) in enumerate(cards):
        x = x0 + i * (cw + 0.25)
        rect(s, x, 1.3, cw, 5.9, WHITE, line=True)
        rect(s, x, 1.3, cw, 0.08, c)
        txt(s, t, x + 0.28, 1.75, cw - 0.5, 1.3, 20, c, bold=True, line_spacing=1.15)
        txt(s, d, x + 0.28, 3.25, cw - 0.5, 3.6, 15.5, DARK, line_spacing=1.3)
    rect(s, 0.83, 7.55, 16.1, 1.5, LIGHT)
    txt(s, "The practical consequence for this project:  predicting where a tumour will spread is not "
           "one prediction problem.\nIt has to be asked separately for each destination.",
        1.2, 7.55, 15.4, 1.5, 18, DARK, bold=True, anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.3)
    footer(s)
    notes(s,
          "Three consequences. The first is conceptual: if you are building a signature of metastasis "
          "and you pool liver, lung and nodal deposits together, you are averaging across a real "
          "difference and the shared part is only about half of it. The second is the one I would most "
          "value your opinion on. In the treatment naive cohort the liver deposits were immune poor "
          "and the lung deposits immune rich, and that matches what is described clinically about "
          "liver metastases responding badly to immunotherapy. The third is about our own tooling. A "
          "model that reads expression from a slide does not transfer between organs, so accuracy "
          "quoted across pooled organs is optimistic.")


# ================================================================== ASK
def ask_slide():
    s = slide(); header(s, "Where your input would help most")
    cards = [("01", "Confirm the B cells on tissue",
              "The clearest difference between liver and nodal deposits was B cell content. A stain "
              "on a handful of matched deposits would confirm or kill it quickly, and does not need "
              "any sequencing.", BLUE),
             ("02", "More patients with two sites",
              "Everything paired rests on four people who gave tissue from two different sites. Even "
              "three or four more such patients would change how strongly this can be stated.",
              ORANGE),
             ("03", "Outcome data if it exists",
              "We have no survival or treatment response for these cases. Without it we cannot ask "
              "whether the site specific part of the program matters clinically.", GREEN)]
    cw, x0 = 5.2, 0.83
    for i, (n, t, d, c) in enumerate(cards):
        x = x0 + i * (cw + 0.25)
        rect(s, x, 1.3, cw, 5.7, WHITE, line=True)
        rect(s, x, 1.3, cw, 0.08, c)
        rect(s, x + 0.28, 1.62, 1.0, 0.5, c)
        txt(s, "ASK " + n, x + 0.28, 1.62, 1.0, 0.5, 12, WHITE, bold=True,
            align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        txt(s, t, x + 0.28, 2.4, cw - 0.5, 1.2, 20, c, bold=True, line_spacing=1.15)
        txt(s, d, x + 0.28, 3.35, cw - 0.5, 3.2, 15.5, DARK, line_spacing=1.3)
    rect(s, 0.83, 7.35, 16.1, 1.6, LIGHT)
    txt(s, "The first one is cheap and would settle the most important open question.  "
           "It needs a pathologist's eye, not more computation.",
        1.2, 7.35, 15.4, 1.6, 18, DARK, bold=True, anchor=MSO_ANCHOR.MIDDLE)
    footer(s)
    notes(s,
          "Three asks, in order of how much they would help. The first is the one I care about most "
          "and it is the cheapest. If B cells really are consistently higher in nodal deposits than "
          "liver deposits from the same patient, that should be visible on a stain, and it does not "
          "need any sequencing to check. If you can point me at a few matched cases, or tell me the "
          "result is already known and obvious, either answer saves us months. The second is more "
          "patients with two metastatic sites. The third is outcome data, which we currently do not "
          "have at all.")


# ================================================================== CLOSING
def closing_slide():
    s = slide()
    rect(s, 0, 0, W, H, BLUE)
    txt(s, "Thank you", 1.5, 3.0, 15, 1.2, 48, WHITE, bold=True)
    rect(s, 1.55, 4.4, 2.6, 0.05, ORANGE)
    txt(s, "Where a pancreatic tumour spreads changes what it becomes.\n"
           "Liver and lymph node deposits share only about half of their departure from the primary "
           "tumour,\nand the part that differs is immune.  The pattern held in a second, independent "
           "group of patients.\n\n"
           "The question I would most like your view on:  are B cells visibly higher in nodal "
           "deposits than in liver deposits?",
        1.5, 4.9, 15.2, 3.0, 19, LIGHT, line_spacing=1.35)
    txt(s, "Spatial Transcriptomics Research   |   14 August 2026", 1.5, 8.7, 12, 0.5, 15, LIGHT)
    footer(s)
    notes(s,
          "Happy to take questions, and I have every number behind these slides if you want to "
          "interrogate any of them. The single most useful thing you could tell me is whether the B "
          "cell difference between nodal and liver deposits is something you would expect to see on a "
          "stain, and whether it is already well known, because that determines whether this is a "
          "finding or a rediscovery.")


# ================================================================== BUILD
title_slide()
overview_slide()

divider("01", "What we measure and who it came from", "The question",
        "Spatial transcriptomics gives us a stained image and the gene readout at\n"
        "the same point in the same tissue.",
        "Before any result, two slides on what the data actually is, because everything afterwards "
        "depends on it. The unusual thing about this technique is that it keeps the image and the "
        "molecular readout registered to the same physical location.")

fig_slide("One slide, measured three ways at every point",
          os.path.join(SS, "s1_what_we_measure.png"),
          "This is the only kind of data that keeps the microscope image and the molecular readout locked to the same physical location.",
          "This is the whole method in one picture. The tissue section is covered with a grid of "
          "measuring points, each 55 microns across, which is about one to ten cells. At every point "
          "we get three things. The stained image of that point, which is the view you would read "
          "down a microscope. The gene readout, about eighteen thousand genes measured at once. And "
          "an estimate of the cell mixture at that point, so what fraction is tumour, immune, stroma, "
          "hepatocyte and so on. The important part is that all three describe the same physical "
          "location, which is what lets us ask whether the image alone predicts the molecular state.")

cohorts_slide()

divider("02", "The main comparison", "Is a liver deposit the same\nas a lymph node deposit?",
        "Four patients gave us tumour from both sites, so the comparison happens\n"
        "inside one person and patient differences cancel out.",
        "This is the core of the talk. The comparison people usually make is metastasis against "
        "primary tumour. We are making a different one, metastasis against metastasis, in the same "
        "person. That design matters because it removes patient to patient variation entirely.")

fig_slide("Only about half of what changes is shared between the two destinations",
          os.path.join(SS, "s2_half_shared.png"),
          "Measured across 42 features of cell content and tissue programme, in the patients who gave tissue from both sites.",
          "Here is the main number. We took everything that changes when a tumour goes from the "
          "pancreas to a metastatic site, forty two different measures of cell content and tissue "
          "programme. We did that separately for the liver deposits and the lymph node deposits, then "
          "asked how much the two changes agree. The answer is about fifty five percent. So a bit "
          "over half of becoming a metastasis is the same wherever it goes, and just under half "
          "depends on the destination. That is the finding. If metastasis were a single switch being "
          "thrown, this number would be close to one hundred percent.")

fig_slide("The difference between the two sites is immune",
          os.path.join(SS, "s2b_bcells.png"),
          "With only four paired patients no p value can reach significance, so what carries the weight is that the direction never reverses.",
          "Having established that the two sites differ, the next question is what differs. We tested "
          "forty two features one at a time and corrected for the fact that we were testing forty "
          "two. Only one immune feature survived, and it is B cells. The left panel is the one I "
          "would look at. Each line is a single patient measured at both of their own sites, and "
          "every single one goes up from liver to lymph node. With four patients I cannot give you a "
          "conventional p value, so what I am leaning on is that the direction never reverses. The "
          "obvious objection is that a lymph node is a lymphoid organ, so of course it has more B "
          "cells. That is exactly what the next slide addresses.")

fig_slide("Could this just be normal lymph node tissue contaminating the sample?",
          os.path.join(SS, "s2c_purity.png"),
          "This is the single control that decides whether the previous slide means anything.  The effect behaves like biology, not like contamination.",
          "This is the control that decides whether the previous slide means anything. The worry is "
          "that we are not measuring tumour biology at all, we are measuring how much surrounding "
          "normal organ got included in the sample. So we repeated the entire analysis while "
          "progressively demanding purer and purer tumour at each measuring point. If the effect were "
          "contamination it would shrink away as purity rose. What you see is that contamination "
          "genuinely does fall, liver cell content drops about eightfold, but the immune difference "
          "between the two sites stays at roughly the same level throughout. That is the pattern you "
          "would expect if the difference is real, and not the pattern you would expect from "
          "contamination.")

divider("03", "An independent check", "Does it happen in other patients?",
        "A separate published cohort, from another centre, with no chemotherapy\n"
        "before sampling, and three metastatic sites instead of two.",
        "A single small cohort is not enough, particularly with a result this surprising. So we took "
        "somebody else's published data, from a different centre, and ran the same analysis. Two "
        "things make it a stronger test than ours: nobody in it had chemotherapy before sampling, and "
        "it covers three different metastatic sites.")

fig_slide("In a separate group of patients, liver deposits are immune poor and lung deposits immune rich",
          os.path.join(SS, "s3_immune_by_site.png"),
          "A different centre, different patients, and no chemotherapy before sampling.  The same pattern appears without us going looking for it.",
          "This is the replication, and I find it the most striking picture in the talk. Each group of "
          "three bars is one immune cell type. Orange is liver deposits, magenta is lung, purple is "
          "peritoneum. Look at the pattern rather than any individual bar. Every single immune "
          "population falls in the liver and rises in the lung. Not some of them, all of them. And "
          "remember none of these patients had chemotherapy before the sample was taken, so this is "
          "not a treatment effect. The clinical reading is that the liver is behaving as an immune "
          "tolerant site and the lung is not, which I believe matches what is seen with "
          "immunotherapy response. I would value your view on whether that is already established.")

fig_slide("A sanity check on the same measurement",
          os.path.join(SS, "s3b_sanity_check.png"),
          "One control did fail: mesothelin cannot mark peritoneal lining here, because pancreatic tumour cells express it themselves.  We report that as a failure.",
          "A quick but important check. Before believing what this measurement says about immune "
          "cells, we should confirm it says the obvious things correctly. It does. It finds liver "
          "cells in the liver deposits and lung tissue in the lung deposits, both clearly. I will "
          "mention one control that failed, because it is instructive. We also tried to detect "
          "peritoneal lining tissue using mesothelin, and it did not work, because pancreatic tumour "
          "cells express mesothelin themselves. So that particular control is uninformative in this "
          "disease and we report it as a failure rather than quietly dropping it.")

divider("04", "Consequences", "What this means for reading slides",
        "If tissue context governs gene expression, a model trained on one organ\n"
        "should fail on another. We tested that directly.",
        "The last piece of evidence is different in kind. The first two results are observations. This "
        "one is a prediction we made from those observations and then tested, which is a stronger form "
        "of argument, and it is in tissue that is not pancreas.")

fig_slide("A model that reads genes from a picture does not carry across organs",
          os.path.join(SS, "s4_cross_organ.png"),
          "This was predicted from the two earlier results and then tested, in six cancers that are not pancreatic.  It is the strongest of the three arguments.",
          "We trained a model to predict gene expression from the stained image and then asked how "
          "well it carries to a new setting. The obvious objection is batch effect, that a model "
          "failing on a new organ is really just failing on a new laboratory with a different scanner "
          "and protocol. So the middle bar holds the organ fixed and changes only the laboratory. "
          "That costs fifteen percent. Changing the organ costs a further thirty seven percent on top. "
          "So the loss is genuinely about tissue, not about technical differences between centres. "
          "One thing worth knowing: this analysis initially gave the opposite answer, because three "
          "sections in the collection were labelled as whole transcriptome but were actually a "
          "targeted panel of about a thousand genes. Filtering those out reversed the conclusion. The "
          "control decided the result, not the headline number.")

negatives_slide()
meaning_slide()
ask_slide()
closing_slide()

prs.save(DEST)
print("wrote", DEST, f"({os.path.getsize(DEST)/1024/1024:.1f} MB, {len(prs.slides.__iter__.__self__._sldIdLst)} slides)")
